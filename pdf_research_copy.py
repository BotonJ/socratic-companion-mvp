#!/usr/bin/env python3
"""
文档文本提取 + Chunk分割

支持格式：
  - PDF（调用pdftotext）
  - PPT/PPTX（使用 python-pptx）
  - TXT 文件
  - 剪贴板（支持 macOS/Linux/Windows）

用法：
  # 处理单个文件
  python3 doc_extract.py input.pdf > output.md
  python3 doc_extract.py input.pptx > output.md
  python3 doc_extract.py input.txt > output.md

  # 从剪贴板读取
  python3 doc_extract.py > output.md

  # 自定义chunk大小
  python3 doc_extract.py input.pdf --chunk-size 3000 > output.md

  # 添加frontmatter
  python3 doc_extract.py input.pdf --frontmatter title="标题" author="作者" > output.md

  # 批量处理目录（PPT专用）
  python3 doc_extract.py --batch /path/to/ppt/dir --output ./materials/

  # 批量处理，保留目录结构
  python3 doc_extract.py --batch /path/to/ppt/dir --output ./materials/ --preserve-structure
"""

import re
import sys
import os
import subprocess
import argparse
import platform
from typing import List, Optional
from datetime import datetime
from pathlib import Path

# 尝试导入 python-pptx
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def is_pdf_file(filepath: str) -> bool:
    """判断是否为PDF文件"""
    return filepath.lower().endswith('.pdf')


def is_ppt_file(filepath: str) -> bool:
    """判断是否为PPT文件"""
    return filepath.lower().endswith(('.ppt', '.pptx'))


def pdf_to_text(pdf_path: str) -> str:
    """
    调用 pdftotext 提取PDF文本
    支持多种编码检测
    """
    # 检查 pdftotext 是否可用
    try:
        subprocess.run(['pdftotext', '-v'],
                       stdout=subprocess.DEVNULL,
                       stderr=subprocess.DEVNULL,
                       check=True)
    except (subprocess.CalledProcessError, FileNotFoundError):
        print("错误：pdftotext 未安装", file=sys.stderr)
        print("macOS: brew install poppler", file=sys.stderr)
        print("Ubuntu/Debian: sudo apt-get install poppler-utils", file=sys.stderr)
        print("Windows: scoop install poppler 或 choco install poppler", file=sys.stderr)
        sys.exit(1)

    try:
        with open(pdf_path, 'rb') as f:
            pdf_data = f.read()
        result = subprocess.run(
            ['pdftotext', '-layout', '-', '-'],
            input=pdf_data,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            check=True
        )
        raw = result.stdout

        # 尝试多种编码
        encodings = ['utf-8', 'latin-1', 'gbk', 'iso-8859-1', 'cp1252']
        for encoding in encodings:
            try:
                text = raw.decode(encoding)
                if '\ufffd' in text[:1000] and encoding == 'utf-8':
                    continue
                return text
            except UnicodeDecodeError:
                continue

        return raw.decode('utf-8', errors='replace')

    except subprocess.CalledProcessError as e:
        print(f"错误：PDF提取失败 - {e}", file=sys.stderr)
        sys.exit(1)


def ppt_to_text(ppt_path: str) -> str:
    """
    使用 python-pptx 提取PPT文本
    返回格式化的Markdown文本
    """
    if not HAS_PPTX:
        print("错误：python-pptx 未安装", file=sys.stderr)
        print("请运行: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    try:
        prs = Presentation(ppt_path)
        lines = []

        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # 跳过重复的页眉（通常很短且以"项目"开头）
                    if len(text) < 15 and (text.startswith("项目") or text.startswith("任务")):
                        continue
                    slide_texts.append(text)

            if slide_texts:
                lines.append(f"## 第 {i} 页")
                lines.append("")
                for text in slide_texts:
                    lines.append(text)
                    lines.append("")

        return '\n'.join(lines)

    except Exception as e:
        print(f"错误：PPT提取失败 - {e}", file=sys.stderr)
        sys.exit(1)


def ppt_to_md(ppt_path: str, title: str = None) -> str:
    """
    将PPT转换为结构化的Markdown文件
    适合教学使用
    """
    if not HAS_PPTX:
        print("错误：python-pptx 未安装", file=sys.stderr)
        print("请运行: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    try:
        prs = Presentation(ppt_path)
        lines = []

        # 从文件名提取标题
        if title is None:
            title = Path(ppt_path).stem

        lines.append(f"# {title}")
        lines.append("")
        lines.append(f"> 来源：PPT课件")
        lines.append(f"> 提取时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}")
        lines.append("")

        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            title_text = None

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()

                    # 识别标题（通常在前面，较短）
                    if not title_text and len(text) < 50 and not text.startswith("（"):
                        # 检查是否像标题（不以标点结尾）
                        if not text[-1] in '。！？，、；：':
                            title_text = text
                            continue

                    # 跳过页眉
                    if len(text) < 15 and (text.startswith("项目") or text.startswith("任务")):
                        continue

                    slide_texts.append(text)

            if slide_texts or title_text:
                lines.append(f"## 第 {i} 页")
                lines.append("")

                if title_text:
                    lines.append(f"### {title_text}")
                    lines.append("")

                for text in slide_texts:
                    # 处理列表项
                    if text.startswith("（") or text.startswith("(") or re.match(r'^\d+[\.、）]', text):
                        lines.append(f"- {text}")
                    else:
                        lines.append(text)
                    lines.append("")

        return '\n'.join(lines)

    except Exception as e:
        print(f"错误：PPT转换失败 - {e}", file=sys.stderr)
        sys.exit(1)


def read_text_from_file(filepath: str) -> str:
    """读取文本文件"""
    with open(filepath, encoding='utf-8') as f:
        return f.read()


def read_text_from_clipboard() -> str:
    """从剪贴板读取文本（跨平台）"""
    system = platform.system()

    try:
        if system == 'Darwin':  # macOS
            return subprocess.check_output(['pbpaste']).decode('utf-8')
        elif system == 'Linux':
            return subprocess.check_output(['xclip', '-o', '-selection', 'clipboard']).decode('utf-8')
        elif system == 'Windows':
            # 方法1：pyperclip（推荐）
            try:
                import pyperclip
                return pyperclip.paste()
            except ImportError:
                # 方法2：PowerShell
                result = subprocess.run(
                    ['powershell', '-command', 'Get-Clipboard'],
                    capture_output=True, text=True
                )
                return result.stdout
        else:
            print(f"错误：不支持的系统 {system}", file=sys.stderr)
            sys.exit(1)
    except Exception as e:
        print(f"错误：无法从剪贴板读取 - {e}", file=sys.stderr)
        sys.exit(1)


def clean_soft_hyphen(text: str) -> str:
    """清理 soft hyphen (\xad) 字符"""
    return text.replace('\xad', '')


def remove_page_numbers(line: str) -> str:
    """移除页码"""
    line = re.sub(r'^\s*\d+\s*', '', line)
    line = re.sub(r'^\s*-\s*\d+\s*-\s*', '', line)
    line = re.sub(r'\s*\d+\s*$', '', line)
    return line.strip()


def normalize_whitespace(text: str) -> str:
    """标准化空白字符"""
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text


def extract_paragraphs(text: str, keep_page_numbers: bool = False) -> List[str]:
    """提取段落"""
    text = clean_soft_hyphen(text)
    text = normalize_whitespace(text)

    lines = text.split('\n')
    paragraphs = []
    current_para = []

    for line in lines:
        if not keep_page_numbers:
            line = remove_page_numbers(line)
        else:
            line = line.strip()

        if not line:
            if current_para:
                paragraph = ' '.join(current_para)
                if paragraph:
                    paragraphs.append(paragraph)
                current_para = []
            continue

        current_para.append(line)

    if current_para:
        paragraph = ' '.join(current_para)
        if paragraph:
            paragraphs.append(paragraph)

    return paragraphs


def split_by_char_count(text: str, max_chars: int) -> List[str]:
    """按字符数分割（在单词边界切断）"""
    chunks = []
    current_chunk = []
    current_length = 0

    break_chars = '.!?,;。！？，；:：'
    word_breaks = ' \t\n'

    for char in text:
        current_chunk.append(char)
        current_length += 1

        if current_length >= max_chars - 50:
            if char in break_chars:
                chunks.append(''.join(current_chunk).strip())
                current_chunk = []
                current_length = 0
            elif char in word_breaks and current_length > max_chars:
                chunks.append(''.join(current_chunk[:-1]).strip())
                current_chunk = []
                current_length = 0

    if current_chunk:
        chunks.append(''.join(current_chunk).strip())

    return [c for c in chunks if c]


def split_by_sentence(text: str, max_chars: int) -> List[str]:
    """按句子边界分割"""
    sentences = re.split(r'([.!?。！？]\s+)', text)
    chunks = []
    current_chunk = []

    for part in sentences:
        if not current_chunk:
            current_chunk = [part]
            continue

        test_chunk = ''.join(current_chunk) + part
        if len(test_chunk) <= max_chars:
            current_chunk.append(part)
        else:
            if current_chunk:
                chunks.append(''.join(current_chunk).strip())
            current_chunk = [part]

    if current_chunk:
        chunks.append(''.join(current_chunk).strip())

    return chunks


def split_by_paragraphs(paragraphs: List[str], max_paragraphs: int) -> List[str]:
    """按段落数量分割"""
    chunks = []
    current_chunk = []

    for para in paragraphs:
        current_chunk.append(para)
        if len(current_chunk) >= max_paragraphs:
            chunks.append('\n\n'.join(current_chunk))
            current_chunk = []

    if current_chunk:
        chunks.append('\n\n'.join(current_chunk))

    return chunks


def generate_frontmatter(frontmatter_dict: dict, chunk_count: int, total_chars: int) -> str:
    """生成YAML frontmatter"""
    lines = ['---']

    for key, value in frontmatter_dict.items():
        if isinstance(value, bool):
            lines.append(f"{key}: {str(value).lower()}")
        elif isinstance(value, str) and '\n' in value:
            lines.append(f"{key}: |")
            for line in value.split('\n'):
                lines.append(f"  {line}")
        else:
            lines.append(f"{key}: {value}")

    if 'date' not in frontmatter_dict:
        lines.append(f"date: {datetime.now().strftime('%Y-%m-%d')}")
    if 'time' not in frontmatter_dict:
        lines.append(f"time: {datetime.now().strftime('%H:%M:%S')}")
    if 'chunk_count' not in frontmatter_dict:
        lines.append(f"chunk_count: {chunk_count}")
    if 'total_chars' not in frontmatter_dict:
        lines.append(f"total_chars: {total_chars}")

    lines.append('---')
    lines.append('')
    return '\n'.join(lines)


def format_chunk(chunk: str, index: int, total: int,
                chunk_prefix: str = '', chunk_suffix: str = '',
                chunk_template: str = None) -> str:
    """格式化单个chunk"""
    if chunk_template:
        chars = len(chunk)
        return chunk_template.format(index=index, total=total, chunk=chunk, chars=chars)
    else:
        parts = []
        if chunk_prefix:
            prefix = chunk_prefix.replace('{i}', str(index)).replace('{n}', str(total))
            parts.append(prefix)
        parts.append(chunk)
        if chunk_suffix:
            suffix = chunk_suffix.replace('{i}', str(index)).replace('{n}', str(total))
            parts.append(suffix)
        return '\n'.join(parts)


def process_single_file(filepath: str, args) -> tuple:
    """
    处理单个文件
    返回 (chunks, paragraphs, total_chars)
    """
    if is_pdf_file(filepath):
        print(f"正在提取PDF: {filepath}", file=sys.stderr)
        text = pdf_to_text(filepath)
    elif is_ppt_file(filepath):
        print(f"正在提取PPT: {filepath}", file=sys.stderr)
        text = ppt_to_text(filepath)
    else:
        print(f"正在读取文件: {filepath}", file=sys.stderr)
        text = read_text_from_file(filepath)

    paragraphs = extract_paragraphs(text, keep_page_numbers=args.no_page_numbers)

    if args.max_paragraphs > 0:
        chunks = split_by_paragraphs(paragraphs, args.max_paragraphs)
    elif args.chunk_size > 0:
        full_text = '\n\n'.join(paragraphs)
        if args.split_by_sentence:
            chunks = split_by_sentence(full_text, args.chunk_size)
        else:
            chunks = split_by_char_count(full_text, args.chunk_size)
    else:
        chunks = ['\n\n'.join(paragraphs)]

    total_chars = sum(len(c) for c in chunks)

    return chunks, paragraphs, total_chars


def batch_process_ppt(input_dir: str, output_dir: str, preserve_structure: bool = False):
    """
    批量处理PPT目录

    Args:
        input_dir: 输入目录路径
        output_dir: 输出目录路径
        preserve_structure: 是否保留原目录结构
    """
    if not HAS_PPTX:
        print("错误：python-pptx 未安装", file=sys.stderr)
        print("请运行: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    input_path = Path(input_dir)
    output_path = Path(output_dir)

    if not input_path.exists():
        print(f"错误：输入目录不存在 - {input_dir}", file=sys.stderr)
        sys.exit(1)

    # 创建输出目录
    output_path.mkdir(parents=True, exist_ok=True)

    # 查找所有PPT文件
    ppt_files = list(input_path.rglob("*.pptx")) + list(input_path.rglob("*.ppt"))

    if not ppt_files:
        print(f"警告：未找到PPT文件 - {input_dir}", file=sys.stderr)
        return

    print(f"找到 {len(ppt_files)} 个PPT文件", file=sys.stderr)

    success_count = 0
    fail_count = 0

    for ppt_file in ppt_files:
        try:
            # 计算输出路径
            if preserve_structure:
                rel_path = ppt_file.relative_to(input_path)
                output_file = output_path / rel_path.with_suffix('.md')
                output_file.parent.mkdir(parents=True, exist_ok=True)
            else:
                output_file = output_path / (ppt_file.stem + '.md')

            # 转换PPT
            md_content = ppt_to_md(str(ppt_file))

            # 写入文件
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(md_content)

            print(f"✅ {ppt_file.name} → {output_file.name}", file=sys.stderr)
            success_count += 1

        except Exception as e:
            print(f"❌ {ppt_file.name}: {e}", file=sys.stderr)
            fail_count += 1

    print(f"\n---", file=sys.stderr)
    print(f"批量处理完成:", file=sys.stderr)
    print(f"  成功: {success_count}", file=sys.stderr)
    print(f"  失败: {fail_count}", file=sys.stderr)
    print(f"  输出目录: {output_path}", file=sys.stderr)


def main():
    parser = argparse.ArgumentParser(
        description='文档文本提取 + Chunk分割（支持PDF/PPT/TXT）',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  # 处理PDF文件
  python3 doc_extract.py input.pdf > output.md

  # 处理PPT文件
  python3 doc_extract.py input.pptx > output.md

  # 处理TXT文件
  python3 doc_extract.py input.txt > output.md

  # 从剪贴板
  python3 doc_extract.py > output.md

  # 自定义chunk大小
  python3 doc_extract.py input.pdf --chunk-size 3000 > output.md

  # 添加frontmatter
  python3 doc_extract.py input.pdf --frontmatter title="标题" type="paper" > output.md

  # 批量处理PPT目录
  python3 doc_extract.py --batch /path/to/ppts --output ./materials/

  # 批量处理，保留目录结构
  python3 doc_extract.py --batch /path/to/ppts --output ./materials/ --preserve-structure
        """
    )

    parser.add_argument('input_file', nargs='?',
                       help='输入文件（PDF/PPT/TXT，可选，默认从剪贴板）')
    parser.add_argument('--chunk-size', type=int, default=2000,
                       help='按字符数分割chunk（默认：2000）')
    parser.add_argument('--max-paragraphs', type=int, default=0,
                       help='按段落数量分割chunk（0=不分割）')
    parser.add_argument('--split-by-sentence', action='store_true',
                       help='按句子边界分割（需要--chunk-size）')
    parser.add_argument('--chunk-prefix', type=str,
                       help='chunk前缀，可用{i}=索引,{n}=总数')
    parser.add_argument('--chunk-suffix', type=str,
                       help='chunk后缀，可用{i}=索引,{n}=总数')
    parser.add_argument('--chunk-template', type=str,
                       help='自定义chunk模板')
    parser.add_argument('--no-page-numbers', action='store_true',
                       help='保留页码')
    parser.add_argument('--frontmatter', nargs='*', metavar='KEY=VALUE',
                       help='添加YAML frontmatter')

    # 批量处理选项
    parser.add_argument('--batch', type=str, metavar='DIR',
                       help='批量处理PPT目录')
    parser.add_argument('--output', type=str, metavar='DIR',
                       help='批量处理输出目录')
    parser.add_argument('--preserve-structure', action='store_true',
                       help='批量处理时保留目录结构')

    args = parser.parse_args()

    # 批量处理模式
    if args.batch:
        output_dir = args.output or './materials/'
        batch_process_ppt(args.batch, output_dir, args.preserve_structure)
        return

    # 解析frontmatter
    frontmatter_dict = {}
    if args.frontmatter:
        for item in args.frontmatter:
            if '=' in item:
                key, value = item.split('=', 1)
                value = value.strip('"\'')
                frontmatter_dict[key] = value
            else:
                frontmatter_dict[item] = True

    # 读取文本
    if args.input_file:
        chunks, paragraphs, total_chars = process_single_file(args.input_file, args)
    else:
        text = read_text_from_clipboard()
        paragraphs = extract_paragraphs(text, keep_page_numbers=args.no_page_numbers)
        chunks = ['\n\n'.join(paragraphs)]
        total_chars = sum(len(c) for c in chunks)

    # 输出frontmatter
    if frontmatter_dict:
        fm = generate_frontmatter(frontmatter_dict, len(chunks), total_chars)
        print(fm)

    # 输出chunks
    for i, chunk in enumerate(chunks, 1):
        formatted = format_chunk(
            chunk, i, len(chunks),
            args.chunk_prefix or '',
            args.chunk_suffix or '',
            args.chunk_template
        )
        print(formatted)
        if i < len(chunks):
            print()

    # 统计信息
    print(f"\n---", file=sys.stderr)
    print(f"提取统计:", file=sys.stderr)
    print(f"  段落数: {len(paragraphs)}", file=sys.stderr)
    print(f"  Chunk数: {len(chunks)}", file=sys.stderr)
    print(f"  总字符: {total_chars}", file=sys.stderr)
    print(f"  平均chunk: {total_chars // len(chunks) if chunks else 0} 字符", file=sys.stderr)


if __name__ == '__main__':
    main()